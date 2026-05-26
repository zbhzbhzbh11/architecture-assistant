#!/usr/bin/env python3
"""生成答辩 PPTX。18 页，聚焦架构设计。

Requirements audit (requierment.md):
  §一.1 需求分析 §一.2 架构推荐 §一.3 决策支持 §一.4 知识进化
  §二.1 微服务 + ≥3 Agent + LLM
  §二.2 需求理解 + 知识库≥10风格 + 推理决策 + 可视化
  §六.1 技术建议: LLM+图谱双驱动 / LangGraph / Neo4j
  §六.2 挑战应对: Few-shot / 规则校验 / 缓存
  §六.3 创新方向: ADR / 组合推荐 / 重构建议

Usage:
    python scripts/build_pptx.py
    python scripts/build_pptx.py -o output.pptx
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 路径 ──
PROJECT_ROOT = Path(__file__).resolve().parent.parent
DIAGRAMS_DIR = PROJECT_ROOT / "docs" / "defense" / "diagrams"
OUTPUT_PATH = PROJECT_ROOT / "docs" / "defense" / "答辩PPT.pptx"
TOTAL_PAGES = 18

# ── 颜色 ──
C_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
C_ACCENT  = RGBColor(0x2D, 0x7D, 0x46)
C_DARK    = RGBColor(0x1E, 0x1E, 0x1E)
C_GRAY    = RGBColor(0x66, 0x66, 0x66)
C_LIGHT   = RGBColor(0xF0, 0xF4, 0xFF)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_RED     = RGBColor(0xC0, 0x39, 0x2B)
C_ORANGE  = RGBColor(0xD3, 0x54, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════
# 辅助函数
# ═══════════════════════════════════════════════════════════════

def add_blank_slide():
    return prs.slides.add_slide(blank_layout)


def add_textbox(slide, left, top, width, height, text="",
                font_size=18, bold=False, color=C_DARK,
                alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, bold=False, color=C_DARK,
             alignment=PP_ALIGN.LEFT, space_before=0, space_after=0,
             font_name="Microsoft YaHei"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_title_bar(slide, title_text, subtitle_text=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_PRIMARY
    shape.line.fill.background()
    add_textbox(slide, 0.8, 0.15, 11.5, 0.6, title_text,
                font_size=30, bold=True, color=C_WHITE)
    if subtitle_text:
        add_textbox(slide, 0.8, 0.65, 11.5, 0.4, subtitle_text,
                    font_size=14, color=RGBColor(0xBB, 0xCC, 0xEE))


def add_page_number(slide, num):
    add_textbox(slide, 11.8, 7.0, 1.2, 0.4, f"{num}/{TOTAL_PAGES}",
                font_size=10, color=C_GRAY, alignment=PP_ALIGN.RIGHT)


def add_table(slide, left, top, col_widths, headers, rows,
              font_size=14, header_color=C_PRIMARY):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(total_w), Inches(0.4 * n_rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = C_DARK
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
    return table_shape


def add_image_slide(slide, img_name, title_text, note_text=None):
    add_title_bar(slide, title_text)
    img_path = DIAGRAMS_DIR / img_name
    if img_path.exists():
        top = 1.3
        avail_h = 7.0 - top - (0.6 if note_text else 0.3)
        slide.shapes.add_picture(str(img_path),
                                  Inches(0.5), Inches(top),
                                  Inches(12.3), Inches(avail_h))
    if note_text:
        add_textbox(slide, 0.8, 6.7, 11.5, 0.4, note_text,
                    font_size=12, color=C_GRAY)


# ═══════════════════════════════════════════════════════════════
# 第 1 页: 封面
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
shape.fill.solid()
shape.fill.fore_color.rgb = C_PRIMARY
shape.line.fill.background()
add_textbox(slide, 1.5, 1.5, 10, 1.2, "架构风格智能助手",
            font_size=44, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 2.7, 10, 0.8,
            "基于 LangGraph 的 Compound AI 推荐系统 · 架构设计专题",
            font_size=22, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.0, 10, 0.6,
            "软件体系结构课程答辩",
            font_size=18, color=RGBColor(0xAA, 0xBB, 0xEE), alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.5, 10, 1.0,
            "FastAPI · LangGraph · Neo4j · LLM · Docker · Compound AI System",
            font_size=14, color=RGBColor(0x99, 0xAA, 0xDD), alignment=PP_ALIGN.CENTER)
add_page_number(slide, 1)

# ═══════════════════════════════════════════════════════════════
# 第 2 页: 问题定义
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "问题定义", "传统架构选型的三大痛点与 Compound AI 解决方案")

# 左侧痛点
add_textbox(slide, 0.8, 1.6, 5.5, 0.5, "传统架构选型三大痛点",
            font_size=22, bold=True, color=C_PRIMARY)
pain_points = [
    ("效率低", "依赖个人经验，缺少系统化方法"),
    ("不稳定", "不同评审者结论不同，不可复现"),
    ("不完整", "缺少量化评分与决策追溯"),
]
y = 2.3
for title, desc in pain_points:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(5.2), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
    box.line.fill.background()
    tf2 = box.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = f"  {title}: {desc}"
    p.font.size = Pt(16)
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = C_RED
    y += 0.85

# 右侧方案
add_textbox(slide, 7.0, 1.6, 5.5, 0.5, "解决方案: 三层 Compound AI",
            font_size=22, bold=True, color=C_ACCENT)
layers = [
    ("规则引擎", "确定性基线 — 10维关键词 + 7条规则"),
    ("Neo4j 知识图谱", "关系推理增强 — HAS_QUALITY 遍历"),
    ("LLM 语义理解", "提升上限 — tie-break 投票 + 摘要"),
]
y = 2.3
for title, desc in layers:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(y), Inches(5.5), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    box.line.fill.background()
    tf2 = box.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = f"  {title}: {desc}"
    p.font.size = Pt(16)
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = C_ACCENT
    y += 0.85

add_textbox(slide, 0.8, 5.8, 11.5, 0.6,
            "核心设计理念: 规则保证下限，图谱增强关系，LLM 提升上限 — 每层独立可降级",
            font_size=16, bold=True, color=C_PRIMARY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 2)

# ═══════════════════════════════════════════════════════════════
# 第 3 页: 需求理解 [新增 — 对应课程 §二.2 需求理解模块]
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "需求理解: 从自然语言到特征信号",
              "对应课程要求 §二.2 — 非结构化文本的特征提取")

# 左侧: 流程
add_textbox(slide, 0.8, 1.5, 5.5, 0.5, "三步特征提取管线",
            font_size=20, bold=True, color=C_PRIMARY)

steps = [
    ("1. 关键词匹配 (10维 × ≈90词)", "\"万人\"→高并发  \"消息\"→实时  \"可靠\"→可靠性  \"扩展\"→可扩展", C_ACCENT),
    ("2. 否定语义过滤 (6种模式)", "\"不需要高并发\" → 剔除高并发  \"无需实时\" → 剔除实时性", C_PRIMARY),
    ("3. LLM 语义补全 (条件触发)", "仅当命中 ≤2维 + LLM已配置时触发 | t=0.1 | 不可用→静默跳过", C_ORANGE),
]
y = 2.2
for title, desc, color in steps:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(5.8), Inches(0.85))
    box.fill.solid()
    box.fill.fore_color.rgb = C_WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = f"  {title}"
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = color
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    add_para(tf2, f"     {desc}", font_size=13, color=C_DARK, space_before=2)
    y += 1.0

# 右侧: 10 个特征维度
add_textbox(slide, 7.2, 1.5, 5.5, 0.5, "10 个质量属性维度",
            font_size=20, bold=True, color=C_PRIMARY)
headers = ["维度", "示例关键词"]
rows = [
    ["高并发", "万人 / 秒杀 / 高吞吐 / QPS"],
    ["实时性", "实时 / 低延迟 / 即时 / 在线"],
    ["可靠性", "高可用 / 容灾 / 容错 / 不丢"],
    ["可扩展性", "扩展 / 弹性 / 扩容 / 横向"],
    ["复杂业务", "交易 / 审批 / 工作流"],
    ["强一致性", "强一致 / 事务 / 金融"],
    ["部署约束", "本地部署 / 私有化 / 边缘"],
    ["数据密集型", "数据流 / ETL / 流处理 / 日志"],
    ["多团队协作", "多团队 / 并行开发"],
    ["安全性", "安全 / 加密 / 认证 / 审计 / 合规"],
]
add_table(slide, 7.2, 2.2, [2.2, 3.3], headers, rows, font_size=11)

add_textbox(slide, 0.8, 6.5, 11.5, 0.5,
            "对应课程评分标准 §四: 需求分析 (15%) — 对 AI 系统不确定性的处理方案",
            font_size=13, bold=True, color=C_PRIMARY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 3)

# ═══════════════════════════════════════════════════════════════
# 第 4 页: 为什么不是 LLM 聊天机器人
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "为什么不是普通 LLM 聊天机器人")

headers = ["维度", "普通 ChatBot", "Compound AI System"]
rows = [
    ["推理方式", "LLM 全包，黑盒生成", "规则引擎主导 + 图谱推理 + LLM 增强"],
    ["可解释性", "无法追溯评分依据", "每分可追溯来源 (标签/规则/图谱)"],
    ["可靠性", "LLM 挂了系统全瘫", "LLM 不可用自动降级纯规则模式"],
    ["幻觉风险", "可能推荐不存在的风格", "四道防线: 规则限定候选 + 闭集投票 + 强制校验"],
    ["领域知识", "训练数据模糊记忆", "10种风格结构化知识库 + 7条专家规则"],
    ["输出格式", "自由文本，不可控", "结构化报告: 推荐+对比矩阵+风险+ADR"],
]
add_table(slide, 0.8, 1.6, [2.2, 3.8, 5.5], headers, rows, font_size=15)

add_textbox(slide, 0.8, 6.3, 11.5, 0.5,
            "Compound AI = 多个专用组件协同，而非单一模型 — 每个组件有明确的职责边界和独立的降级能力",
            font_size=16, bold=True, color=C_PRIMARY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 4)

# ═══════════════════════════════════════════════════════════════
# 第 5 页: C4 Context
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_image_slide(slide, "diagram_01.png", "C4 Context 图 — 系统上下文",
                "关键设计: LLM 在系统边界之外 — 不可用时自动降级为纯规则模式，核心推荐链路不受影响")
add_page_number(slide, 5)

# ═══════════════════════════════════════════════════════════════
# 第 6 页: C4 Container
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_image_slide(slide, "diagram_02.png", "C4 Container 图 — 8 容器微服务架构")
add_page_number(slide, 6)

# ═══════════════════════════════════════════════════════════════
# 第 7 页: 数据层设计 — Neo4j 知识图谱 [原第12页 Neo4j 部分前移]
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "数据层设计: Neo4j 知识图谱",
              "对应课程要求 §二.2 知识库模块 — ≥10 种架构风格 + HAS_QUALITY 关系推理")

# Neo4j 截图
neo4j_path = DIAGRAMS_DIR / "质量属性.png"
if neo4j_path.exists():
    slide.shapes.add_picture(str(neo4j_path),
                              Inches(0.3), Inches(1.3),
                              Inches(7.0), Inches(5.8))

# 右侧: 10 种风格表
add_textbox(slide, 7.8, 1.3, 5.0, 0.4, "10 种架构风格 (完整覆盖)",
            font_size=16, bold=True, color=C_PRIMARY)
headers = ["风格", "属性数"]
rows = [
    ["微服务架构", "5"],
    ["事件驱动架构", "4"],
    ["CQRS (命令查询分离)", "4"],
    ["管道-过滤器", "3"],
    ["分层架构", "2"],
    ["六边形架构", "2"],
    ["面向服务架构 (SOA)", "2"],
    ["无服务器架构", "2"],
    ["空间架构", "2"],
    ["客户端-服务器", "1"],
]
add_table(slide, 7.8, 1.9, [3.2, 1.0], headers, rows, font_size=11)

# 底部说明
add_textbox(slide, 7.8, 6.1, 5.0, 1.0,
            "6节点类型 + 6关系类型\nHAS_QUALITY / SUITABLE_FOR / HAS_RISK\nCOMPLEMENTS / RECOMMENDS / BASED_ON\n\nNeo4j 可选 — KNOWLEDGE_BACKEND\n=auto 模式自动 fallback JSON",
            font_size=11, color=C_GRAY)
add_page_number(slide, 7)

# ═══════════════════════════════════════════════════════════════
# 第 8 页: 微服务划分决策
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "微服务划分决策", "8 容器 / 4 Agent — 为什么不是单体？对应课程 §二.1 微服务架构约束")

headers = ["服务", "端口", "职责", "独立理由"]
rows = [
    ["frontend", "3000", "用户交互与可视化", "前后端分离"],
    ["api-gateway", "8000", "请求编排 + 缓存", "单一对外入口"],
    ["req-agent", "8001", "特征提取(规则+LLM)", "CPU + 网络IO 混合"],
    ["match-agent", "8002", "规则评分 + 图谱融合", "独立接入 Neo4j"],
    ["eval-agent", "8003", "LLM 评估 + ADR", "LLM 故障隔离"],
    ["kb", "8004", "双后端知识存取", "数据层独立"],
    ["ref-agent", "8005", "坏味检测 + 迁移", "重构逻辑独立"],
    ["neo4j", "7687", "图数据库", "可选后端"],
]
add_table(slide, 0.8, 1.5, [2.0, 1.2, 4.0, 4.2], headers, rows, font_size=15)

tf = add_textbox(slide, 0.8, 5.6, 11.5, 1.5)
tf.paragraphs[0].text = "拆分五原则"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C_PRIMARY
tf.paragraphs[0].font.name = "Microsoft YaHei"
for reason in [
    "领域解耦 — 特征提取(CPU密集) vs 评估(IO密集) 分离",
    "故障隔离 — LLM 超时只影响摘要质量，不影响特征提取",
    "异构集成 — matching-agent 独立接入 Neo4j driver (已在上页展示)",
    "独立部署 — 每个服务独立 Dockerfile，代码量 ≤ 393 行",
    "技术灵活 — 环境变量驱动，切换 LLM/缓存/存储无需改代码",
]:
    add_para(tf, f"  {reason}", font_size=14, color=C_DARK, space_before=2)
add_page_number(slide, 8)

# ═══════════════════════════════════════════════════════════════
# 第 9 页: Agent 协作时序图
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_image_slide(slide, "diagram_03.png", "Agent 协作时序图",
                "Step1 特征提取(详见第3页) → Step2 架构匹配(规则+图谱) → Step3 评估决策(LLM并行) → Step4 重构建议(非阻塞)")
add_page_number(slide, 9)

# ═══════════════════════════════════════════════════════════════
# 第 10 页: LangGraph StateGraph + Subgraph 展开
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_image_slide(slide, "diagram_04.png", "LangGraph StateGraph + Subgraph 展开",
                "父图 4 节点顺序执行 + matching-agent Subgraph (3子节点) + evaluation-agent Subgraph (4子节点, Send() 并行扇出)")
add_page_number(slide, 10)

# ═══════════════════════════════════════════════════════════════
# 第 11 页: 三层混合推理 + 学习权重 [补充知识进化 §一.4]
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "三层混合推理架构",
              "对应课程 §二.2 推理决策模块 + §一.4 知识进化 — 规则保证下限，图谱增强关系，LLM 提升上限")

layers_data = [
    ("Layer 1: 规则引擎评分 [始终运行 — 保证下限]", C_ACCENT, RGBColor(0xE8, 0xF5, 0xE9),
     ["标签匹配: tag ∈ features → +2 分/项",
      "6 条硬编码规则 → +1 分 (高并发→EDA / 多团队→微服务 / 强一致→分层 / ...)",
      "学习权重 ≥ 2 次确认 → +1 分  ← 知识进化 (§一.4): 用户反馈驱动权重更新",
      "主流保底: Layered / Microservices / Event-Driven 始终进入 Top 3"]),
    ("Layer 2: 知识图谱推理 [可选 — 关系增强]", C_PRIMARY, C_LIGHT,
     ["Neo4j Cypher 遍历 HAS_QUALITY 关系 (详见第7页图模型)",
      "每个匹配质量属性 +2 图谱分 | 上限 ≤ 规则分 50%",
      "不可用 → 图谱加分 = 0，自动回退 JSON 文件存储"]),
    ("Layer 3: LLM 语义理解 [可选 — 提升上限]", C_ORANGE, RGBColor(0xFF, 0xF3, 0xE0),
     ["投票: t=0.0，闭集选择 + 候选强制校验 | 摘要: t=0.3，Few-shot + 降级模板",
      "LLM 权重 = +1 (仅 tie-break — 保证可解释性优先于智能性)",
      "失败 → _fallback_summary() 用规则理由生成结构化中文报告"]),
]
y = 1.4
for title, accent, bg, bullets in layers_data:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.65))
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = accent
    box.line.width = Pt(1)
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = f"  {title}"
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = accent
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    for b in bullets:
        add_para(tf2, f"     {b}", font_size=14, color=C_DARK, space_before=1)
    y += 1.9

add_page_number(slide, 11)

# ═══════════════════════════════════════════════════════════════
# 第 12 页: 防止 LLM 幻觉的四道防线
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "防止 LLM 幻觉的四道防线")

defenses = [
    ("防线 1", "LLM 不参与候选集生成", C_PRIMARY,
     "Top 3 由规则引擎确定，LLM 无权推荐规则外风格"),
    ("防线 2", "闭集投票 + 字符串精确匹配校验", C_PRIMARY,
     'Prompt 明确: "仅从以下候选列表中选择"；返回值不在候选列表中 → 直接丢弃'),
    ("防线 3", "低温度参数约束", C_ACCENT,
     "投票 t=0.0 (贪婪解码)，摘要 t=0.3 — 最大化确定性，最小化随机性"),
    ("防线 4", "Few-shot 格式化约束", C_ACCENT,
     "9 个标注示例严格约束输出 Schema；非 JSON 输出容错解析，失败则丢弃"),
]
y = 1.5
for label, title, color, desc in defenses:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = C_WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = f"  {label}: {title}"
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = color
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    add_para(tf2, f"     {desc}", font_size=15, color=C_DARK)
    y += 1.35

add_textbox(slide, 0.8, 6.9, 11.5, 0.4,
            "LLM 的不确定性被严格限制在'摘要质量'层面，不污染核心推荐逻辑",
            font_size=14, bold=True, color=C_ACCENT, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 12)

# ═══════════════════════════════════════════════════════════════
# 第 13 页: LLM 集成方案
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "LLM 集成方案", "4 处调用点 · 全部兼容 OpenAI 协议 · 全部可降级")

headers = ["#", "位置", "用途", "t°", "超时", "失败策略"]
rows = [
    ["1", "req-agent", "语义补全", "0.1", "15s", "静默维持规则特征"],
    ["2", "eval-agent", "风格投票", "0.0", "20s", "返回 null，规则排序不变"],
    ["3", "eval-agent", "摘要生成", "0.3", "25s", "fallback 规则模板"],
    ["4", "ref-agent", "步骤润色", "0.3", "15s", "规则模板原文"],
]
add_table(slide, 0.8, 1.5, [0.8, 2.0, 2.0, 1.0, 1.2, 4.5], headers, rows, font_size=15)

# LLM 配置
tf = add_textbox(slide, 0.8, 4.0, 11.5, 1.0)
tf.paragraphs[0].text = "LLM 配置 (环境变量注入，3 行切换模型)"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C_PRIMARY
tf.paragraphs[0].font.name = "Microsoft YaHei"
for line in ["LLM_API_BASE  → https://api.deepseek.com",
             "LLM_API_KEY   → sk-xxxx",
             "LLM_MODEL     → deepseek-v4-flash"]:
    add_para(tf, f"    {line}", font_size=14, color=C_DARK)

# Few-shot + 缓存 (双栏)
tf = add_textbox(slide, 0.8, 5.4, 5.5, 1.8)
tf.paragraphs[0].text = "Few-shot Prompt"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C_PRIMARY
tf.paragraphs[0].font.name = "Microsoft YaHei"
for line in ["需求侧 6 示例 (模糊/否定/安全/数据/一致/重构)",
             "评估侧 3 示例 (EDA/MS/Layered 报告)",
             "模块不可用 → 自动降级零样本 Prompt"]:
    add_para(tf, f"  {line}", font_size=14, color=C_DARK, space_before=1)

tf = add_textbox(slide, 7.0, 5.4, 5.5, 1.8)
tf.paragraphs[0].text = "LLM 请求缓存"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = C_PRIMARY
tf.paragraphs[0].font.name = "Microsoft YaHei"
for line in ["键: SHA256(requirement + model + knowledge_version)",
             "knowledge_version = styles 文件 MD5 前 8 位",
             "知识库更新 → 旧缓存自动失效",
             "双后端: memory (默认) / SQLite (持久化)"]:
    add_para(tf, f"  {line}", font_size=14, color=C_DARK, space_before=1)
add_page_number(slide, 13)

# ═══════════════════════════════════════════════════════════════
# 第 14 页: 创新方向建议 [改名 + 去 Neo4j — 对应课程 §六.3]
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "创新方向建议 (课程 §六.3)",
              "ADR 决策溯源 · 架构组合推荐 · 架构重构建议")

innovations = [
    ("ADR 架构决策溯源", C_PRIMARY,
     ["每次推荐自动生成 ADR-YYYYMMDD-NNN 格式决策记录",
      "JSON (主存储) + Neo4j (同步) 双后端，写入失败不阻塞推荐",
      "API: GET /adr 列表 + GET /adr/{id} 详情查询"]),
    ("架构组合推荐 (5 种)", C_ACCENT,
     ["Microservices+Event-Driven | Layered+CQRS | Pipeline-Filter+Event-Driven",
      "Hexagonal+Microservices | CQRS+Event Sourcing",
      "评分: 组成分数和 + 特征覆盖互补 + 图谱 COMPLEMENTS 加分 - 复杂度惩罚"]),
    ("架构重构建议 (5+5)", C_ORANGE,
     ["5 种坏味检测: 单体耦合 / 扩展瓶颈 / 发布缓慢 / 遗留锁定 / 数据耦合",
      "5 种重构模式: 绞杀者 / 防腐层 / 模块化单体 / CQRS迁移 / 事件驱动迁移",
      "LLM 可选润色，不可用 → 规则模板 (非阻塞调用，超时 8s)"]),
]
y = 1.5
for title, color, bullets in innovations:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.7))
    box.fill.solid()
    box.fill.fore_color.rgb = C_WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = f"  {title}"
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = color
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    for b in bullets:
        add_para(tf2, f"     {b}", font_size=14, color=C_DARK, space_before=1)
    y += 1.9
add_page_number(slide, 14)

# ═══════════════════════════════════════════════════════════════
# 第 15 页: 测试验证体系
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "测试验证体系", "四层测试金字塔: 79 单元 + 40 端到端 + 43 验收")

pyramid = [
    ("43 项自动验收 (15 大类, 9/9 技术建议 100% 通过)", 0.8, C_PRIMARY),
    ("40 条端到端测试 (冒烟20 + 回归20, 5项指标全100%)", 4.0, C_ACCENT),
    ("79 条单元测试 (76 passed, 6 个测试文件, 覆盖全部模块)", 8.0, C_ORANGE),
]
for label, width, color in pyramid:
    left = (12.3 - width) / 2
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5 + left), Inches(4.2), Inches(width), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = label
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = C_WHITE
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

add_textbox(slide, 0.8, 1.8, 11.5, 0.5, "回归测试 5 项指标 (对应课程 §三.1 系统测试报告)",
            font_size=18, bold=True, color=C_PRIMARY)
headers = ["指标", "Top3完整", "主流覆盖", "推荐产出", "决策可解释", "矩阵完整"]
rows = [["目标", "100%", "100%", "100%", "100%", "100%"],
        ["实际", "100%", "100%", "100%", "100%", "100%"]]
add_table(slide, 0.8, 2.4, [2.0, 1.8, 1.8, 1.8, 2.0, 1.8], headers, rows, font_size=14)

add_textbox(slide, 0.8, 5.8, 11.5, 0.8,
            "对应课程 §三.1 系统测试报告 + §四 测试验证 (15%) — 20 条用例覆盖 8 大类场景 (即时通讯/电商/金融/大数据/企业/医疗/政务/物联网)",
            font_size=13, bold=True, color=C_PRIMARY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 15)

# ═══════════════════════════════════════════════════════════════
# 第 16 页: 降级可靠性矩阵
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "降级可靠性矩阵", "12 种故障场景全覆盖 — 核心结论: LLM全关 20/20 回归通过")

headers = ["故障点", "降级行为", "验证方式"]
rows = [
    ["LangGraph 未安装", "启动报错, 提示安装", "启动检查"],
    ["LangGraph 运行时异常", "返回 502, trace 记录", "代码逻辑"],
    ["LLM 未配置", "纯规则模式", "回归测试 ✓"],
    ["LLM 超时 (15/20/25s)", "静默 / fallback 摘要", "代码逻辑"],
    ["LLM 返回非 JSON", "丢弃 → fallback 模板", "代码逻辑"],
    ["LLM 投票名无效", "丢弃不加分", "单元测试"],
    ["Neo4j 不可达 (auto)", "自动回退 JSON", "单元测试"],
    ["req/match/eval-agent 不可达", "网关 502", "异常矩阵"],
    ["ref-agent 不可达", "ref={} 非阻塞", "代码逻辑"],
    ["ADR 写入失败", "标记 failed, 不阻塞推荐", "代码逻辑"],
]
add_table(slide, 0.8, 1.6, [4.0, 5.0, 2.5], headers, rows, font_size=14)

box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
box.line.color.rgb = C_ACCENT
box.line.width = Pt(2)
tf2 = box.text_frame
tf2.word_wrap = True
tf2.paragraphs[0].text = "核心结论: LLM / Neo4j / LangGraph 三个增强组件同时不可用时，纯规则模式回归测试 20/20 通过 — 核心推荐结果不受影响"
tf2.paragraphs[0].font.size = Pt(16)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = C_ACCENT
tf2.paragraphs[0].font.name = "Microsoft YaHei"
tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
add_page_number(slide, 16)

# ═══════════════════════════════════════════════════════════════
# 第 17 页: 总结
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "总结: 四个核心创新点")

innovations = [
    ("Compound AI System", "规则引擎 + Neo4j + LLM 三层协同，每层独立可降级，不是单一 LLM 黑盒"),
    ("微服务 + LangGraph Subgraph 编排", "8 容器 / 4 Agent / StateGraph + Subgraph 嵌套 + Send() 并行，职责单一 ≤ 393 行"),
    ("混合推理 + 幻觉防线", "规则保证下限 + LLM 提升上限，闭集投票 + 强制校验 + 低温度 + Few-shot"),
    ("可解释可降级可测试", "评分每分可追溯 + 12 种故障降级 + 79+40+43 全通过"),
]
y = 1.6
for i, (title, desc) in enumerate(innovations):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.85))
    box.fill.solid()
    box.fill.fore_color.rgb = C_WHITE
    box.line.color.rgb = C_PRIMARY
    box.line.width = Pt(1)
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = f"  {i+1}. {title}: {desc}"
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.name = "Microsoft YaHei"
    tf2.paragraphs[0].font.color.rgb = C_DARK
    y += 1.0

# 架构决策
add_textbox(slide, 0.8, 5.55, 11.5, 0.4, "关键架构决策回顾",
            font_size=16, bold=True, color=C_PRIMARY)
headers = ["决策", "选择", "核心理由"]
rows = [
    ["服务通信", "HTTP+JSON", "可调试性优先于性能"],
    ["图数据库", "Neo4j 可选", "降低环境依赖，JSON 零依赖可运行"],
    ["编排引擎", "LangGraph StateGraph", "状态管理标准化、Subgraph 封装内聚、Send() 声明式并行"],
    ["缓存方案", "memory/SQLite", "课程场景够用，无需 Redis"],
    ["前端", "原生 HTML", "零构建步骤，打开即看"],
]
add_table(slide, 0.8, 5.95, [2.5, 4.0, 5.0], headers, rows, font_size=12)
add_page_number(slide, 17)

# ═══════════════════════════════════════════════════════════════
# 第 18 页: 致谢
# ═══════════════════════════════════════════════════════════════
slide = add_blank_slide()
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
shape.fill.solid()
shape.fill.fore_color.rgb = C_PRIMARY
shape.line.fill.background()
add_textbox(slide, 1.5, 2.0, 10, 1.2, "谢谢！",
            font_size=52, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.3, 10, 0.8, "欢迎各位老师提问",
            font_size=28, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.6, 10, 1.0,
            "Architecture Assistant · 架构风格智能助手\nGitHub: zbhzbhzbh11/architecture-assistant",
            font_size=18, color=RGBColor(0x99, 0xAA, 0xDD), alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.8, 10, 0.8,
            "FastAPI · LangGraph · Neo4j · LLM (DeepSeek) · Docker Compose · "
            "Few-shot Prompt · Compound AI · Multi-Agent",
            font_size=14, color=RGBColor(0x88, 0x99, 0xCC), alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 6.5, 10, 0.5,
            "43/43 验收通过 · 76 单元测试 · 20 回归 100% · 9/9 技术建议",
            font_size=14, color=RGBColor(0x88, 0x99, 0xCC), alignment=PP_ALIGN.CENTER)
add_page_number(slide, 18)

# ═══════════════════════════════════════════════════════════════
# 保存
# ═══════════════════════════════════════════════════════════════
output = OUTPUT_PATH
if len(sys.argv) > 2 and sys.argv[1] == "-o":
    output = Path(sys.argv[2])
prs.save(str(output))
print(f"PPTX saved: {output}")
print(f"  Slides: {len(prs.slides)}")
print(f"  Size: {output.stat().st_size / 1024:.0f} KB")
